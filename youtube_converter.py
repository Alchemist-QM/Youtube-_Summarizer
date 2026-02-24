import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

#MAX_SLIDES = 20


class YoutubeConverter:
    def __init__(self,file_path:str="final_output.json"):
        self.MAX_SLIDES = 20
        self.file_path = file_path
        
    def open_json_file(self):
        try:
            with open(self.file_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except:
            raise FileNotFoundError

        
    def create_pptx_from_data(self,):
        data = self.open_json_file()
        timeline = data["timeline"][:self.MAX_SLIDES]

        prs = Presentation()

        # --- Intro Slide ---
        intro_slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = intro_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = "Introduction:\n" + data.get("introduction", "") + "\n\nTLDR:\n" + data.get("tldr", "")
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.LEFT

        # --- Main Slides ---
        for entry in timeline:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            img_path = entry["frame"]
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(1), height=Inches(3))
            txBox = slide.shapes.add_textbox(Inches(4.5), Inches(1), Inches(4.5), Inches(3))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = "Summary: " + entry.get("summary", "")
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.LEFT
            # Optionally add transcript
            if entry.get("transcript"):
                tf.add_paragraph().text = "Transcript: " + " ".join(entry["transcript"])

        # --- Outro Slide ---
        outro_slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = outro_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = (
            "Final Summary:\n" + data.get("final_summary", "") +
            "\n\nFinal Insights:\n" + data.get("final_insights", "") +
            "\n\nConclusion:\n" + data.get("conclusion", "")
        )
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.LEFT

        prs.save("slideshow.pptx")
        print("PowerPoint saved as slideshow.pptx")






    def create_html_from_data(self,):
        data = self.open_json_file()
        timeline = data["timeline"][:self.MAX_SLIDES]

        html_head = """
        <!DOCTYPE html>
        <html>
        <head>
        <title>Video Slideshow</title>
        <style>
        .slide { display: flex; margin-bottom: 40px; }
        .slide img { width: 400px; margin-right: 40px; }
        .slide .text { font-size: 1.2em; }
        .intro, .outro { margin-bottom: 50px; font-size: 1.3em; background: #eef; padding: 2em; border-radius: 10px; }
        </style>
        </head>
        <body>
        <h1>Video Slideshow</h1>
        """

        html_tail = """
        </body>
        </html>
        """

        # --- Intro Slide ---
        intro_html = f'''
        <div class="intro">
            <h2>Introduction</h2>
            <p>{data.get("introduction", "")}</p>
            <h3>TLDR</h3>
            <p>{data.get("tldr", "")}</p>
        </div>
        '''

        # --- Main Slides ---
        slides_html = ""
        for idx, entry in enumerate(timeline):
            slides_html += f"""
        <div class="slide">
            <img src="{entry['frame']}" alt="Frame {idx}">
            <div class="text">
                <b>Summary:</b> {entry.get('summary', '')}<br>
                <b>Transcript:</b> {'<br>'.join(entry.get('transcript', []))}
            </div>
        </div>
        """

        # --- Outro Slide ---
        outro_html = f'''
        <div class="outro">
            <h2>Final Summary</h2>
            <p>{data.get("final_summary", "")}</p>
            <h3>Final Insights</h3>
            <p>{data.get("final_insights", "")}</p>
            <h3>Conclusion</h3>
            <p>{data.get("conclusion", "")}</p>
        </div>
        '''

        with open("slideshow.html", "w", encoding="utf-8") as f:
            f.write(html_head + intro_html + slides_html + outro_html + html_tail)

        print("HTML Slideshow saved as slideshow.html")